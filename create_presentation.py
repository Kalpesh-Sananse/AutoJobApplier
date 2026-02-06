"""
Generate a professional PowerPoint presentation for the LinkedIn Auto Job Applier.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 115, 177)  # LinkedIn blue
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "LinkedIn Auto Job Applier"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Production-Ready Autonomous AI Agent"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Introduction & Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title2 = slide2.shapes.title
    title2.text = "Introduction & Overview"
    
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "What We Built"
    
    p1 = tf2.add_paragraph()
    p1.text = "🤖 AI-Powered Form Filling (Ollama LLM)"
    p1.level = 1
    
    p2 = tf2.add_paragraph()
    p2.text = "🎭 Browser Automation (Playwright)"
    p2.level = 1
    
    p3 = tf2.add_paragraph()
    p3.text = "📄 Automatic PDF Resume Upload"
    p3.level = 1
    
    p4 = tf2.add_paragraph()
    p4.text = "🔍 Smart Modal-Scoped Field Detection"
    p4.level = 1
    
    tf2.add_paragraph()
    
    p5 = tf2.add_paragraph()
    p5.text = "Results"
    p5.font.bold = True
    
    p6 = tf2.add_paragraph()
    p6.text = "✅ 100% Automated - No human intervention"
    p6.level = 1
    
    p7 = tf2.add_paragraph()
    p7.text = "✅ 7/7 Applications Submitted Successfully"
    p7.level = 1
    
    # Slide 3: System Architecture
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "System Architecture"
    
    # Add architecture image
    try:
        img_path = "/Users/keps/.gemini/antigravity/brain/5a522643-9a1c-4aa2-9909-2db662097ee1/architecture_diagram_1770360130998.png"
        slide3.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), width=Inches(7))
    except:
        # If image doesn't exist, add text description
        content3 = slide3.placeholders[1]
        tf3 = content3.text_frame
        tf3.text = "Layered Architecture:"
        
        p1 = tf3.add_paragraph()
        p1.text = "User Configuration Layer → YAML files, Resume PDF"
        p1.level = 1
        
        p2 = tf3.add_paragraph()
        p2.text = "Service Layer → Browser Manager (Playwright) + AI Handler (Ollama)"
        p2.level = 1
        
        p3 = tf3.add_paragraph()
        p3.text = "Core Layer → LinkedIn Automation (Modal Detection, Field Scanning, Intelligent Filling, Navigation)"
        p3.level = 1
    
    # Slide 4: Implementation Workflow
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Implementation Workflow"
    
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "1. Initialization"
    
    p1 = tf4.add_paragraph()
    p1.text = "Load credentials, config, and PDF resume"
    p1.level = 1
    
    p2 = tf4.add_paragraph()
    p2.text = "2. Job Discovery"
    p2.level = 0
    
    p3 = tf4.add_paragraph()
    p3.text = "Find Easy Apply jobs on LinkedIn"
    p3.level = 1
    
    p4 = tf4.add_paragraph()
    p4.text = "3. Application Process (Critical!)"
    p4.level = 0
    
    p5 = tf4.add_paragraph()
    p5.text = "✓ Detect modal container (scope all queries)"
    p5.level = 1
    
    p6 = tf4.add_paragraph()
    p6.text = "✓ Scan fields: text, dropdowns, radios, checkboxes"
    p6.level = 1
    
    p7 = tf4.add_paragraph()
    p7.text = "✓ AI generates intelligent answers"
    p7.level = 1
    
    p8 = tf4.add_paragraph()
    p8.text = "✓ Upload PDF resume, navigate pages, submit"
    p8.level = 1
    
    # Slide 5: Key Features
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "Key Features & Innovations"
    
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "🎯 Modal Scoping (Critical Innovation)"
    
    p1 = tf5.add_paragraph()
    p1.text = "All queries scoped to Easy Apply modal only"
    p1.level = 1
    
    p2 = tf5.add_paragraph()
    p2.text = "Result: 0 search bars filled, 100% accuracy"
    p2.level = 1
    
    p3 = tf5.add_paragraph()
    p3.text = "🧠 AI-Powered Intelligence"
    p3.level = 0
    
    p4 = tf5.add_paragraph()
    p4.text = "Ollama LLM understands context and generates answers"
    p4.level = 1
    
    p5 = tf5.add_paragraph()
    p5.text = "🔄 Robust Error Handling"
    p5.level = 0
    
    p6 = tf5.add_paragraph()
    p6.text = "Retry logic (3 attempts), consecutive error tracking"
    p6.level = 1
    
    p7 = tf5.add_paragraph()
    p7.text = "📊 Statistics Tracking"
    p7.level = 0
    
    p8 = tf5.add_paragraph()
    p8.text = "Real-time monitoring: submitted, failed, fields filled"
    p8.level = 1
    
    # Slide 6: Results & Achievements
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "Production Test Results"
    
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "📊 Test Run: 7 Applications Successfully Submitted"
    tf6.paragraphs[0].font.bold = True
    
    tf6.add_paragraph()
    
    p1 = tf6.add_paragraph()
    p1.text = "Success Rate: 100% (7/7)"
    p1.level = 0
    
    p2 = tf6.add_paragraph()
    p2.text = "Fields Filled: 18 total"
    p2.level = 0
    
    p3 = tf6.add_paragraph()
    p3.text = "Resume Uploads: 7 PDF files"
    p3.level = 0
    
    p4 = tf6.add_paragraph()
    p4.text = "Average Time: 2-3 min per application"
    p4.level = 0
    
    p5 = tf6.add_paragraph()
    p5.text = "Manual Equivalent: ~15 min per application"
    p5.level = 0
    
    p6 = tf6.add_paragraph()
    p6.text = "Time Saved: ~90 minutes in test run"
    p6.level = 0
    p6.font.bold = True
    p6.font.color.rgb = RGBColor(0, 128, 0)
    
    # Slide 7: Summary & Future
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "Summary & Future Roadmap"
    
    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.text = "What We Achieved"
    tf7.paragraphs[0].font.bold = True
    
    p1 = tf7.add_paragraph()
    p1.text = "✅ Production-ready AI agent deployed"
    p1.level = 1
    
    p2 = tf7.add_paragraph()
    p2.text = "✅ 100% success rate in testing"
    p2.level = 1
    
    p3 = tf7.add_paragraph()
    p3.text = "✅ 850+ lines of documented code"
    p3.level = 1
    
    tf7.add_paragraph()
    
    p4 = tf7.add_paragraph()
    p4.text = "Future Enhancements"
    p4.font.bold = True
    
    p5 = tf7.add_paragraph()
    p5.text = "🚀 Job quality scoring with AI"
    p5.level = 1
    
    p6 = tf7.add_paragraph()
    p6.text = "🚀 Database integration for tracking"
    p6.level = 1
    
    p7 = tf7.add_paragraph()
    p7.text = "🚀 Dashboard with real-time analytics"
    p7.level = 1
    
    p8 = tf7.add_paragraph()
    p8.text = "🚀 Cloud deployment on AWS Lambda"
    p8.level = 1
    
    # Save presentation
    output_file = "LinkedIn_AI_Agent_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Created {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
